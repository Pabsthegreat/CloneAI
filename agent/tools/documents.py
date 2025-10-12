
"""
Document utilities for CloneAI
Robust merging (PDFs, PPTs) and cross‑platform conversions with graceful fallbacks.

Key fixes in this version
- FIX: Correct comtypes usage (use comtypes.CreateObject, not comtypes.client.CreateObject).
- FIX: PDF→DOCX on Windows now uses Microsoft Word COM (no NumPy), avoiding pdf2docx crashes.
- FIX: Output path can be a DIRECTORY; auto-derives filename from input.
- PPT merge fidelity:
    * Windows: PowerPoint COM copy/paste and re-apply source Design.
    * Non-Windows: Convert each PPT(X)→PDF (LibreOffice), then merge PDFs for perfect visuals.
    * Optional python-pptx fallback to produce PPTX (best-effort) if needed.
- Conversions hardened with clear errors & fallbacks (LibreOffice where Office isn’t available).
"""

from __future__ import annotations

import os
import sys
import shutil
import subprocess
import tempfile
from datetime import datetime
from typing import List, Optional, Tuple
from copy import deepcopy  # used only by python-pptx fallback

# Lazy/optional deps
try:
    from PyPDF2 import PdfMerger  # type: ignore
except Exception:
    PdfMerger = None  # type: ignore

Presentation = None  # python-pptx Presentation (lazy)
Converter = None     # pdf2docx Converter (lazy)
comtypes = None      # comtypes (lazy)


# ---------------- helpers ----------------
def _is_windows() -> bool:
    return sys.platform == "win32"


def _ensure_parent_dir(path: str) -> None:
    d = os.path.dirname(os.path.abspath(path))
    if d:
        os.makedirs(d, exist_ok=True)


def _which(cmd: str) -> Optional[str]:
    return shutil.which(cmd)


def _has_libreoffice() -> bool:
    return bool(_which("soffice") or _which("libreoffice"))


def _soffice_bin() -> str:
    return _which("soffice") or _which("libreoffice")  # type: ignore


def _resolve_output_path(input_path: str, output_path: Optional[str], new_ext: str) -> str:
    """
    If output_path is a directory or None, derive filename from input_path and append new_ext.
    If output_path looks like a file (has an extension), return as-is.
    """
    in_base = os.path.splitext(os.path.basename(input_path))[0]
    if not output_path or os.path.isdir(output_path):
        return os.path.join(output_path or os.getcwd(), in_base + new_ext)
    # If no extension given, treat as directory
    root, ext = os.path.splitext(output_path)
    if ext == "":
        return os.path.join(output_path, in_base + new_ext)
    return output_path


def _run_soffice_convert(src_path: str, out_dir: str, fmt: str) -> str:
    """
    Convert file with LibreOffice headless. Returns path to converted file.
    fmt examples: 'pdf', 'pptx', 'docx'
    """
    if not _has_libreoffice():
        raise RuntimeError("LibreOffice not found in PATH. Install LibreOffice to enable this fallback.")

    bin_path = _soffice_bin()
    os.makedirs(out_dir, exist_ok=True)
    cmd = [
        bin_path,
        "--headless",
        "--invisible",
        "--nocrashreport",
        "--nodefault",
        "--view",
        "--nolockcheck",
        "--nologo",
        "--nofirststartwizard",
        "--convert-to",
        fmt,
        "--outdir",
        out_dir,
        os.path.abspath(src_path),
    ]
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        raise RuntimeError(
            f"LibreOffice failed converting '{src_path}' → {fmt}. "
            f"stdout={proc.stdout.strip()} stderr={proc.stderr.strip()}"
        )
    # Determine produced filename
    base = os.path.splitext(os.path.basename(src_path))[0]
    out_path = os.path.join(out_dir, f"{base}.{fmt}")
    if not os.path.exists(out_path):
        # LO sometimes appends suffixes; try to locate
        matches = [f for f in os.listdir(out_dir) if f.startswith(base + ".") and f.endswith("." + fmt)]
        if matches:
            out_path = os.path.join(out_dir, matches[0])
    if not os.path.exists(out_path):
        raise RuntimeError(f"Expected output file not found after conversion to {fmt}: {out_path}")
    return out_path


def _safe_import_presentation():
    global Presentation
    if Presentation is None:
        try:
            from pptx import Presentation as _Presentation  # type: ignore
            Presentation = _Presentation
        except Exception as e:
            raise ImportError("python-pptx is not installed. pip install python-pptx") from e


def _safe_import_pdf2docx():
    global Converter
    if Converter is None:
        # Suppress noisy import warnings (NumPy runtime warnings on some builds)
        import warnings
        warnings.filterwarnings("ignore", category=Warning)
        warnings.filterwarnings("ignore", category=RuntimeWarning)
        try:
            from pdf2docx import Converter as _Converter  # type: ignore
            Converter = _Converter
        except Exception as e:
            raise ImportError("pdf2docx is not installed. pip install pdf2docx") from e


def _safe_import_comtypes():
    global comtypes
    if comtypes is None:
        try:
            import comtypes.client as comtypes  # type: ignore
        except Exception as e:
            raise ImportError("comtypes is not installed. On Windows, pip install comtypes") from e


# ---------------- main API ----------------
class DocumentManager:
    """Manager for document operations (merge, convert)."""

    def __init__(self):
        self.supported_pdf_formats = [".pdf"]
        self.supported_ppt_formats = [".pptx", ".ppt"]  # .ppt will be auto-converted as needed
        self.supported_doc_formats = [".docx", ".doc"]

    # ---- listing ----
    def list_files_in_directory(self, directory: str, file_type: str) -> List[Tuple[str, datetime]]:
        if not os.path.exists(directory):
            raise FileNotFoundError(f"Directory not found: {directory}")

        if file_type == "pdf":
            extensions = self.supported_pdf_formats
        elif file_type == "ppt":
            extensions = self.supported_ppt_formats
        elif file_type == "docx":
            extensions = self.supported_doc_formats
        else:
            raise ValueError(f"Unsupported file type: {file_type}")

        out: List[Tuple[str, datetime]] = []
        for fn in os.listdir(directory):
            p = os.path.join(directory, fn)
            if os.path.isfile(p) and os.path.splitext(fn)[1].lower() in extensions:
                out.append((fn, datetime.fromtimestamp(os.path.getmtime(p))))
        return out

    # ---- PDF merge ----
    def merge_pdfs(
        self,
        directory: str,
        output_path: str,
        file_list: Optional[List[str]] = None,
        start_file: Optional[str] = None,
        end_file: Optional[str] = None,
        order: str = "asc",
    ) -> str:
        if PdfMerger is None:
            raise ImportError("PyPDF2 not installed. pip install PyPDF2")

        # Select files
        if file_list:
            files = [(f, None) for f in file_list]
        elif start_file and end_file:
            all_files = self.list_files_in_directory(directory, "pdf")
            all_files.sort(key=lambda x: x[1])
            names = [n for n, _ in all_files]
            if start_file not in names: raise ValueError(f"Start file not found: {start_file}")
            if end_file not in names:   raise ValueError(f"End file not found: {end_file}")
            i, j = names.index(start_file), names.index(end_file)
            if i > j: i, j = j, i
            files = all_files[i:j+1]
        else:
            files = self.list_files_in_directory(directory, "pdf")
            files.sort(key=lambda x: x[1])

        if order == "desc":
            files.reverse()
        if not files:
            raise ValueError("No PDF files found to merge")

        _ensure_parent_dir(output_path)
        with PdfMerger() as merger:  # type: ignore
            for fname, _ in files:
                merger.append(os.path.join(directory, fname))
            merger.write(output_path)
        return output_path

    # ---- PPT merge ----
    def merge_ppts(
        self,
        directory: str,
        output_path: str,
        file_list: Optional[List[str]] = None,
        start_file: Optional[str] = None,
        end_file: Optional[str] = None,
        order: str = "asc",
        prefer_pdf_on_non_windows: bool = True,
        allow_python_pptx_fallback: bool = True,
    ) -> str:
        """
        Merge multiple PowerPoint files into one.

        Strategy:
          * Windows + Office: COM merge → PPTX (preserves source formatting).
          * Else if output_path endswith '.pdf' OR prefer_pdf_on_non_windows=True:
              Convert each PPT/PPTX → PDF (LibreOffice), then merge PDFs → output.
          * Else if allow_python_pptx_fallback:
              Best-effort PPTX merge via python-pptx (themes/media may drift).
        """
        # Build file list
        if file_list:
            files = [(f, None) for f in file_list]
        elif start_file and end_file:
            all_files = self.list_files_in_directory(directory, "ppt")
            all_files.sort(key=lambda x: x[1])
            names = [n for n, _ in all_files]
            if start_file not in names: raise ValueError(f"Start file not found: {start_file}")
            if end_file not in names:   raise ValueError(f"End file not found: {end_file}")
            i, j = names.index(start_file), names.index(end_file)
            if i > j: i, j = j, i
            files = all_files[i:j+1]
        else:
            files = self.list_files_in_directory(directory, "ppt")
            files.sort(key=lambda x: x[1])

        if order == "desc":
            files.reverse()
        if not files:
            raise ValueError("No PowerPoint files found to merge")

        # Windows COM path
        com_reason = None
        if _is_windows():
            try:
                return self._merge_ppts_via_com(directory, files, output_path)
            except Exception as e:
                com_reason = f"COM merge failed: {e}"

        # PDF path (non-Windows or user explicitly wants .pdf)
        if (not _is_windows() and prefer_pdf_on_non_windows) or output_path.lower().endswith(".pdf"):
            if not output_path.lower().endswith(".pdf"):
                raise ValueError("On non-Windows, set output_path to .pdf for best fidelity or disable prefer_pdf_on_non_windows.")
            with tempfile.TemporaryDirectory() as td:
                pdfs: List[str] = []
                for fname, _ in files:
                    src = os.path.join(directory, fname)
                    pdfs.append(self._ppt_any_to_pdf(src, td))
                merged_tmp = os.path.join(td, "_merged.pdf")
                self._merge_specific_pdfs(pdfs, merged_tmp)
                _ensure_parent_dir(output_path)
                shutil.copyfile(merged_tmp, output_path)
                return output_path

        # python-pptx fallback
        if allow_python_pptx_fallback:
            return self._merge_ppts_via_python_pptx(directory, files, output_path)

        msg = "Unable to merge PowerPoints with perfect fidelity."
        if com_reason:
            msg += " " + com_reason
        raise RuntimeError(msg)

    # ---- conversions ----
    def convert_pdf_to_docx(self, pdf_path: str, output_path: Optional[str]) -> str:
        if not os.path.exists(pdf_path):
            raise FileNotFoundError(f"PDF file not found: {pdf_path}")
        output_path = _resolve_output_path(pdf_path, output_path, ".docx")
        _ensure_parent_dir(output_path)

        # Prefer Word COM on Windows (avoids pdf2docx/numpy issues)
        if _is_windows():
            try:
                _safe_import_comtypes()
                word = comtypes.CreateObject("Word.Application")
                word.Visible = False
                try:
                    # Open PDF and let Word's PDF Reflow convert on SaveAs
                    doc = word.Documents.Open(os.path.abspath(pdf_path), ConfirmConversions=False, ReadOnly=True)
                    try:
                        # 12 = wdFormatXMLDocument (.docx)
                        doc.SaveAs(os.path.abspath(output_path), FileFormat=12)
                    finally:
                        doc.Close(False)
                finally:
                    word.Quit()
                return output_path
            except Exception as e:
                # Fall back below
                pass

        # Next best: LibreOffice
        if _has_libreoffice():
            out_dir = os.path.dirname(os.path.abspath(output_path))
            converted = _run_soffice_convert(pdf_path, out_dir, "docx")
            if os.path.abspath(converted) != os.path.abspath(output_path):
                shutil.move(converted, output_path)
            return output_path

        # Last resort: pdf2docx
        _safe_import_pdf2docx()
        cv = Converter(os.path.abspath(pdf_path))
        try:
            cv.convert(os.path.abspath(output_path))
        finally:
            cv.close()
        return output_path

    def convert_docx_to_pdf(self, docx_path: str, output_path: Optional[str]) -> str:
        if not os.path.exists(docx_path):
            raise FileNotFoundError(f"DOCX file not found: {docx_path}")
        output_path = _resolve_output_path(docx_path, output_path, ".pdf")
        _ensure_parent_dir(output_path)

        if _is_windows():
            _safe_import_comtypes()
            word = comtypes.CreateObject("Word.Application")
            word.Visible = False
            try:
                doc = word.Documents.Open(os.path.abspath(docx_path))
                try:
                    doc.SaveAs(os.path.abspath(output_path), FileFormat=17)  # 17=PDF
                finally:
                    doc.Close(False)
            finally:
                word.Quit()
            return output_path

        # Non-Windows: LibreOffice
        out_dir = os.path.dirname(os.path.abspath(output_path))
        converted = _run_soffice_convert(docx_path, out_dir, "pdf")
        if os.path.abspath(converted) != os.path.abspath(output_path):
            shutil.move(converted, output_path)
        return output_path

    def convert_ppt_to_pdf(self, ppt_path: str, output_path: Optional[str]) -> str:
        if not os.path.exists(ppt_path):
            raise FileNotFoundError(f"PowerPoint file not found: {ppt_path}")
        output_path = _resolve_output_path(ppt_path, output_path, ".pdf")
        _ensure_parent_dir(output_path)

        if _is_windows():
            _safe_import_comtypes()
            app = comtypes.CreateObject("PowerPoint.Application")
            app.Visible = 1
            try:
                pres = app.Presentations.Open(os.path.abspath(ppt_path), WithWindow=False)
                try:
                    pres.SaveAs(os.path.abspath(output_path), 32)  # 32=PDF
                finally:
                    pres.Close()
            finally:
                app.Quit()
            return output_path

        # Non-Windows: LibreOffice
        out_dir = os.path.dirname(os.path.abspath(output_path))
        converted = _run_soffice_convert(ppt_path, out_dir, "pdf")
        if os.path.abspath(converted) != os.path.abspath(output_path):
            shutil.move(converted, output_path)
        return output_path

    # ---- internals ----
    def _merge_specific_pdfs(self, pdf_paths: List[str], output_path: str) -> str:
        if PdfMerger is None:
            raise ImportError("PyPDF2 not installed. pip install PyPDF2")
        _ensure_parent_dir(output_path)
        with PdfMerger() as merger:  # type: ignore
            for p in pdf_paths:
                merger.append(p)
            merger.write(output_path)
        return output_path

    def _ppt_any_to_pdf(self, src: str, out_dir: str) -> str:
        """Convert .ppt or .pptx to PDF using best method for the OS."""
        if _is_windows():
            try:
                return self.convert_ppt_to_pdf(src, os.path.join(out_dir, os.path.splitext(os.path.basename(src))[0] + ".pdf"))
            except Exception:
                if _has_libreoffice():
                    return _run_soffice_convert(src, out_dir, "pdf")
                raise
        else:
            return _run_soffice_convert(src, out_dir, "pdf")

    def _ppt_to_pptx_if_needed(self, src: str, out_dir: str) -> str:
        """Return a .pptx path: if src is .ppt, convert to .pptx; if .pptx, return as-is."""
        ext = os.path.splitext(src)[1].lower()
        if ext == ".pptx":
            return src
        # Convert .ppt → .pptx
        if _is_windows():
            try:
                _safe_import_comtypes()
                app = comtypes.CreateObject("PowerPoint.Application")
                app.Visible = 1
                try:
                    pres = app.Presentations.Open(os.path.abspath(src), WithWindow=False)
                    try:
                        dst = os.path.join(out_dir, os.path.splitext(os.path.basename(src))[0] + ".pptx")
                        pres.SaveAs(os.path.abspath(dst))  # default PPTX
                    finally:
                        pres.Close()
                finally:
                    app.Quit()
                return dst
            except Exception:
                if _has_libreoffice():
                    return _run_soffice_convert(src, out_dir, "pptx")
                raise RuntimeError("Failed to convert .ppt to .pptx; install LibreOffice or Office.")
        else:
            if _has_libreoffice():
                return _run_soffice_convert(src, out_dir, "pptx")
            raise RuntimeError("LibreOffice is required to convert .ppt to .pptx on non-Windows.")

    def _merge_ppts_via_com(self, directory: str, files: List[Tuple[str, datetime]], output_path: str) -> str:
        """Windows-only: merge PPT/PPTX to PPTX using PowerPoint COM, keeping source formatting."""
        _safe_import_comtypes()
        msoTrue, msoFalse = -1, 0
        app = comtypes.CreateObject("PowerPoint.Application")
        app.Visible = msoTrue
        _ensure_parent_dir(output_path)
        out_abs = os.path.abspath(output_path)

        dest = None
        try:
            dest = app.Presentations.Add(msoFalse)
            for fname, _ in files:
                src_path = os.path.abspath(os.path.join(directory, fname))
                src = app.Presentations.Open(src_path, WithWindow=msoFalse)
                try:
                    count = src.Slides.Count
                    if count == 0:
                        continue
                    src.Slides.Range().Copy()
                    dest.Slides.Paste()
                    start_idx = dest.Slides.Count - count + 1
                    for i in range(count):
                        dest_slide = dest.Slides(start_idx + i)
                        src_slide = src.Slides(i + 1)
                        try:
                            dest_slide.Design = src_slide.Design
                        except Exception:
                            pass
                finally:
                    src.Close()
            dest.SaveAs(out_abs)  # PPTX
            return out_abs
        finally:
            if dest is not None:
                pass
            app.Quit()

    def _merge_ppts_via_python_pptx(self, directory: str, files: List[Tuple[str, datetime]], output_path: str) -> str:
        """
        Best-effort PPTX merge without COM (themes/media may drift).
        Safeguards:
          - Auto-convert any .ppt to .pptx via LibreOffice/COM into a temp dir
          - Skip uncopyable shapes rather than crashing
        """
        _safe_import_presentation()
        _ensure_parent_dir(output_path)
        with tempfile.TemporaryDirectory() as td:
            normalized: List[str] = []
            for fname, _ in files:
                src = os.path.join(directory, fname)
                if os.path.splitext(src)[1].lower() == ".pptx":
                    normalized.append(src)
                else:
                    normalized.append(self._ppt_to_pptx_if_needed(src, td))

            merged = Presentation()  # type: ignore
            try:
                blank = merged.slide_layouts[6]
            except IndexError:
                blank = merged.slide_layouts[0]

            for p in normalized:
                src_prs = Presentation(p)  # type: ignore
                for slide in src_prs.slides:
                    new_slide = merged.slides.add_slide(blank)
                    for shape in slide.shapes:
                        try:
                            el = shape.element
                            newel = deepcopy(el)
                            new_slide.shapes._spTree.insert_element_before(newel, "p:extLst")
                        except Exception:
                            # Skip shapes that can't be cloned
                            continue

            if len(merged.slides) == 0:
                merged.slides.add_slide(blank)
            merged.save(output_path)
            return output_path


# --------------- public wrappers ---------------
def merge_pdf_files(
    directory: str,
    output_path: str,
    file_list: Optional[List[str]] = None,
    start_file: Optional[str] = None,
    end_file: Optional[str] = None,
    order: str = "asc",
) -> str:
    return DocumentManager().merge_pdfs(directory, output_path, file_list, start_file, end_file, order)


def merge_ppt_files(
    directory: str,
    output_path: str,
    file_list: Optional[List[str]] = None,
    start_file: Optional[str] = None,
    end_file: Optional[str] = None,
    order: str = "asc",
    prefer_pdf_on_non_windows: bool = True,
    allow_python_pptx_fallback: bool = True,
) -> str:
    return DocumentManager().merge_ppts(
        directory,
        output_path,
        file_list,
        start_file,
        end_file,
        order,
        prefer_pdf_on_non_windows=prefer_pdf_on_non_windows,
        allow_python_pptx_fallback=allow_python_pptx_fallback,
    )


def convert_pdf_to_docx(pdf_path: str, output_path: Optional[str]) -> str:
    return DocumentManager().convert_pdf_to_docx(pdf_path, output_path)


def convert_docx_to_pdf(docx_path: str, output_path: Optional[str]) -> str:
    return DocumentManager().convert_docx_to_pdf(docx_path, output_path)


def convert_ppt_to_pdf(ppt_path: str, output_path: Optional[str]) -> str:
    return DocumentManager().convert_ppt_to_pdf(ppt_path, output_path)


def list_documents_in_directory(directory: str, file_type: str) -> List[Tuple[str, datetime]]:
    return DocumentManager().list_files_in_directory(directory, file_type)
