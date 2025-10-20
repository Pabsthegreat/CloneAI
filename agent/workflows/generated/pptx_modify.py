from agent.workflows import register_workflow
from agent.workflows.registry import ParameterSpec, WorkflowContext
from typing import Dict, Any
import os
from pptx import Presentation
from pptx.util import Pt

@register_workflow(
    namespace="pptx",
    name="modify",
    summary="Auto-generated workflow for `pptx:modify`",
    description="Implements the CLI command `pptx:modify slide:5 message:\"Happy Diwali\" presentation:diwali.pptx`.",
    parameters=(
        ParameterSpec(name="slide", description="The slide number to modify (1-based index)", type=int, required=True),
        ParameterSpec(name="message", description="The message to add to the slide", type=str, required=True),
        ParameterSpec(name="presentation", description="Path to the .pptx file to modify", type=str, required=True),
        ParameterSpec(name="output", description="Optional output path for the modified presentation (defaults to overwrite)", type=str, required=False, default=None),
    ),
    metadata={
        "category": "PPTX COMMANDS",
        "usage": "pptx:modify slide:5 message:\"Happy Diwali\" presentation:diwali.pptx [output:output.pptx]",
        "examples": [
            "pptx:modify slide:5 message:\"Happy Diwali\" presentation:diwali.pptx",
            "pptx:modify slide:2 message:\"Welcome!\" presentation:welcome.pptx output:welcome_updated.pptx"
        ],
    }
)
def pptx_modify_handler(ctx: WorkflowContext, params: Dict[str, Any]) -> str:
    slide_num = params.get("slide")
    message = params.get("message")
    presentation_path = params.get("presentation")
    output_path = params.get("output")

    # Check if file exists, if not check in artifacts folder
    if not os.path.isfile(presentation_path):
        import pathlib
        artifacts_path = pathlib.Path.cwd() / "artifacts" / presentation_path
        if artifacts_path.exists():
            presentation_path = str(artifacts_path)
        else:
            return f"Error: Presentation file '{presentation_path}' does not exist (also checked in artifacts folder)."

    try:
        prs = Presentation(presentation_path)
    except Exception as e:
        return f"Error: Failed to open presentation: {e}"

    if slide_num < 1 or slide_num > len(prs.slides):
        return f"Error: Slide number {slide_num} is out of range. This presentation has {len(prs.slides)} slides."

    slide = prs.slides[slide_num - 1]

    # Add a textbox with the message in the center of the slide
    left = top = width = height = None
    # Default to center of slide
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    width = int(slide_width * 0.6)
    height = int(slide_height * 0.2)
    left = int((slide_width - width) / 2)
    top = int((slide_height - height) / 2)

    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = message
    font = run.font
    font.size = Pt(40)
    font.bold = True
    font.name = 'Arial'
    p.alignment = 1  # Center

    # Save the modified presentation
    save_path = output_path if output_path else presentation_path
    try:
        prs.save(save_path)
    except Exception as e:
        return f"Error: Failed to save presentation: {e}"

    return f"Message '{message}' added to slide {slide_num} in '{save_path}'."
