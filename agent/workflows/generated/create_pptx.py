from agent.workflows import register_workflow
from agent.workflows.registry import ParameterSpec, WorkflowContext
from typing import Dict, Any
from agent.tools.ollama_client import run_ollama
from agent.tools.documents import DocumentManager
import tempfile
import os

@register_workflow(
    namespace="create",
    name="pptx",
    summary="Auto-generated workflow for `create:pptx`",
    description="Implements the CLI command `create:pptx description:5 slides about diwali`.",
    parameters=(
        ParameterSpec(name="description", description="Description of the presentation and number of slides, e.g. '5 slides about diwali'", type=str, required=True),
        ParameterSpec(name="filename", description="Optional output filename for the PPTX file", type=str, required=False, default="presentation.pptx"),
    ),
    metadata={
        "category": "CREATE COMMANDS",
        "usage": "create:pptx description:5 slides about diwali [filename:mydeck.pptx]",
        "examples": ["create:pptx description:5 slides about diwali", "create:pptx description:3 slides about AI filename:ai_intro.pptx"],
    }
)
def create_pptx_handler(ctx: WorkflowContext, params: Dict[str, Any]) -> str:
    description = params.get("description")
    filename = params.get("filename", "presentation.pptx")
    if not description or not isinstance(description, str):
        return "Error: 'description' parameter is required and must be a string."

    # Step 1: Use LLM to generate slide titles and content
    prompt = f"""You are an expert presentation designer. Given the following description, generate a list of slide titles and bullet points for each slide.

Description: {description}

Respond with ONLY valid JSON in this exact format (no markdown, no explanations):
[
  {{"title": "Slide Title 1", "bullets": ["Point 1", "Point 2", "Point 3"]}},
  {{"title": "Slide Title 2", "bullets": ["Point A", "Point B"]}}
]

JSON:"""
    
    from agent.config.runtime import LOCAL_PLANNER
    
    try:
        llm_response = run_ollama(prompt, profile=LOCAL_PLANNER)
    except Exception as e:
        return f"Error: Failed to generate slide content using LLM: {e}"

    try:
        import json
        slides = json.loads(llm_response)
        if not isinstance(slides, list) or not all(isinstance(slide, dict) for slide in slides):
            return "Error: LLM did not return a valid slide structure."
    except Exception as e:
        return f"Error: Failed to parse LLM response as JSON: {e}"

    # Step 2: Generate PPTX using python-pptx
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        return "Error: python-pptx is not installed. Please install it to use this workflow."

    try:
        prs = Presentation()
        for idx, slide_data in enumerate(slides):
            title = slide_data.get("title", f"Slide {idx+1}")
            bullets = slide_data.get("bullets", [])
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = title
            body_shape = slide.shapes.placeholders[1]
            tf = body_shape.text_frame
            tf.clear()
            for bullet in bullets:
                p = tf.add_paragraph()
                p.text = bullet
                p.level = 0
    except Exception as e:
        return f"Error: Failed to generate PPTX: {e}"

    # Step 3: Save PPTX to artifacts folder
    try:
        # Create artifacts folder if it doesn't exist
        import pathlib
        artifacts_dir = pathlib.Path.cwd() / "artifacts"
        artifacts_dir.mkdir(exist_ok=True)
        
        pptx_path = artifacts_dir / filename
        prs.save(str(pptx_path))
    except Exception as e:
        return f"Error: Failed to save PPTX: {e}"

    return f"✅ Presentation created successfully!\n\nSaved to: {pptx_path}\n\n📊 Contains {len(slides)} slides:\n" + "\n".join([f"  {i+1}. {s.get('title', 'Untitled')}" for i, s in enumerate(slides)])
